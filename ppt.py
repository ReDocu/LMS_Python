from pptx import Presentation

# 새 프레젠테이션 생성
prs = Presentation()

# --- 표지 슬라이드 ---
slide_layout = prs.slide_layouts[0]  # 제목 슬라이드
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "자기소개서"
subtitle.text = "개인 이력 및 경력 요약"

# --- 1페이지: 교육, 자격증, 보유기술 ---
slide_layout = prs.slide_layouts[1]  # 제목과 내용
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "교육 & 자격증 & 보유기술"

slide.placeholders[1].text = (
    "■ 교육이수\n"
    "- 2023.09~2024.05 | 국제컴퓨터아트학원 | AI 스마트플랫폼 개발자 양성과정\n"
    "- 2019.09~2020.03 | 경일게임아카데미 | Unity 3D 스마트 콘텐츠 개발 구직자 과정\n\n"
    "■ 자격증\n"
    "- 2016.05 | 정보처리산업기사 | 한국산업인력공단\n"
    "- 2012.11 | 네트워크관리사 2급(국가공인) | 한국정보통신자격협회\n\n"
    "■ 보유기술\n"
    "- C/C++, Unreal, Unity, Python"
)

# --- 2페이지: 경력 & 학력 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "경력 & 학력"

slide.placeholders[1].text = (
    "■ 경력\n"
    "- 2020.07~2021.12 | 셀빅 | 대리 | 2600~2800만원\n"
    "  · Unity/C# AR 콘텐츠 개발, 센서 기반 FPS/퍼즐/낚시 게임 개발\n\n"
    "- 2022.05~2023.05 | KETI | 사원(프리랜서) | 4000~4500만원\n"
    "  · Unreal/Carla 자율주행 시뮬레이션, 플러그인 커스터마이징, UI·환경요소 개발\n\n"
    "■ 학력\n"
    "- 2017.03~2019.08 | 건국대 대학원 | 스마트ICT융합 | 석사 | GPA 4.2/4.5\n"
    "- 2012.03~2017.02 | 학점은행제 | 컴퓨터공학 | 학사 | GPA 3.58/4.5\n"
    "- 2008~2010 | 잠실고등학교 | 졸업"
)

# 파일 저장
prs.save("self_intro_clean.pptx")
print("생성 완료! self_intro_clean.pptx 파일을 확인하세요.")
